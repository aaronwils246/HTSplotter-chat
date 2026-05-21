"""
Generate a 6-slide PowerPoint summary of the HTSplotter analysis.
Run: python make_slides.py
Output: HTSplotter_Summary.pptx

Slide 1 — Title & overview
Slide 2 — The 4 experiment types
Slide 3 — Drug screen results
Slide 4 — Drug combination synergy
Slide 5 — Genetic perturbagen screens
Slide 6 — AI-assisted chat interface
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x37, 0x5E)
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x37, 0x96, 0x37)
ORANGE     = RGBColor(0xC5, 0x5A, 0x11)
TEAL       = RGBColor(0x1F, 0x86, 0x78)
PURPLE     = RGBColor(0x70, 0x30, 0xA0)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
PANEL_DARK = RGBColor(0x0F, 0x23, 0x40)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

TOTAL = 6

# ── Helpers ────────────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=13, bold=False, color=TEXT_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def para(tf, text, size=12, bold=False, color=TEXT_DARK,
         before=5, italic=False, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def header(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 0.88, fill=DARK_BLUE)
    box(slide, 0, 0, 0.08, 0.88, fill=MID_BLUE)
    txt(slide, title, 0.25, 0.1, 12.8, 0.55, size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.25, 0.6, 12.8, 0.25, size=12, color=LIGHT_BLUE)

def footer(slide, n):
    box(slide, 0, 7.28, 13.33, 0.22, fill=DARK_BLUE)
    txt(slide, f"Slide {n} of {TOTAL}", 11.5, 7.3, 1.7, 0.2,
        size=9, color=WHITE, align=PP_ALIGN.RIGHT)
    txt(slide, "HTSplotter  |  Nunes et al., PLoS One 2024  |  DOI: 10.1371/journal.pone.0296322",
        0.25, 7.3, 10.0, 0.2, size=9, color=LIGHT_BLUE)

def bullet_block(slide, items, l, t, w, h, size=11, color=TEXT_DARK, marker="•  "):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(4)
        r = p.add_run(); r.text = marker + item
        r.font.size = Pt(size); r.font.color.rgb = color
    return tb


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title & High-Level Overview
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, DARK_BLUE)
box(s1, 0, 0, 13.33, 0.1, fill=MID_BLUE)

# Title block
txt(s1, "HTSplotter", 0.6, 0.55, 12.1, 1.3,
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, "High-Throughput Screening Analysis",
    0.6, 1.8, 12.1, 0.55, size=24, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
box(s1, 2.0, 2.5, 9.33, 0.05, fill=MID_BLUE)
txt(s1, "All Example Datasets  ·  MCF7 Cells  ·  Confluency Readout  ·  AI Chat Interface",
    0.6, 2.6, 12.1, 0.4, size=15, color=RGBColor(0x88, 0xBB, 0xE0), align=PP_ALIGN.CENTER)

# Left card — what is HTSplotter
box(s1, 0.45, 3.15, 7.6, 3.85, fill=PANEL_DARK)
box(s1, 0.45, 3.15, 0.07, 3.85, fill=MID_BLUE)
tb = s1.shapes.add_textbox(Inches(0.65), Inches(3.25), Inches(7.2), Inches(3.65))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "What is HTSplotter?"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = LIGHT_BLUE

para(tf,
    "HTSplotter is an open-source Python tool for end-to-end processing, "
    "analysis, and visualisation of chemical and genetic in vitro perturbation screens. "
    "It automatically identifies the experiment type, normalises data, fits dose-response "
    "curves, calculates IC values, and scores drug synergy.",
    size=12, color=WHITE, before=10)

para(tf, "Supports four experiment types:", size=12, color=LIGHT_BLUE, before=10)
for item in [
    "Drug screen  —  IC values, dose-response curves, growth rates",
    "Drug combination  —  Bliss / HSA / ZIP synergy scoring",
    "Genetic perturbagen  —  siRNA/shRNA normalised activity",
    "Genetic-chemical  —  gene-knockdown + drug interaction",
]:
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    r2 = p2.add_run(); r2.text = "     ●  " + item
    r2.font.size = Pt(11); r2.font.color.rgb = WHITE

para(tf, "Web tool: htsplotter.cmgg.be  ·  GitHub: github.com/CBIGR/HTSplotter",
    size=10, color=RGBColor(0x88, 0xBB, 0xE0), before=10)

# Right card — citation
box(s1, 8.35, 3.15, 4.55, 3.85, fill=PANEL_DARK)
box(s1, 8.35, 3.15, 0.07, 3.85, fill=MID_BLUE)
tb2 = s1.shapes.add_textbox(Inches(8.55), Inches(3.25), Inches(4.2), Inches(3.65))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]
r2 = p2.add_run(); r2.text = "Citation"
r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = LIGHT_BLUE

para(tf2,
    "Nunes C, Anckaert J, De Vloed F, De Wyn J, Durinck K, "
    "Vandesompele J, Speleman F, Vermeirssen V.",
    size=11, color=WHITE, before=12)
para(tf2,
    "HTSplotter: An end-to-end data processing, analysis and "
    "visualisation tool for chemical and genetic in vitro "
    "perturbation screening.",
    size=11, bold=True, color=WHITE, before=8)
para(tf2, "PLoS One. 2024;19(1):e0296322",
    size=11, italic=True, color=LIGHT_BLUE, before=8)
para(tf2, "DOI: 10.1371/journal.pone.0296322",
    size=11, color=LIGHT_BLUE, before=5)
para(tf2, "PMID: 38181013",
    size=11, color=LIGHT_BLUE, before=4)

footer(s1, 1)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — The 4 Experiment Types
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2, WHITE)
header(s2, "The 4 Experiment Types Supported by HTSplotter",
       "Each type auto-detected from input file header — no manual configuration needed")

TYPES = [
    {
        "title": "Drug Screen",
        "icon":  "💊",
        "color": MID_BLUE,
        "light": RGBColor(0xEF, 0xF4, 0xFB),
        "what":  "Tests a single drug across a concentration range to determine potency.",
        "outputs": ["Dose-response curves (sigmoidal fit)", "IC₂₅ / IC₅₀ / IC₇₅ values",
                    "R² and curve fit statistics", "Growth rate over time (multi-TP)"],
        "datasets": ["1 time point (72 h endpoint)", "37 time points (2 h interval, 74 h total) × 2 variants"],
        "status": "✅  10 / 10 datasets complete",
        "sc": TEAL,
    },
    {
        "title": "Drug Combination",
        "icon":  "🔬",
        "color": PURPLE,
        "light": RGBColor(0xF3, 0xEE, 0xFA),
        "what":  "Tests two drugs in a 7×7 dose matrix to identify synergy or antagonism.",
        "outputs": ["Bliss / HSA / ZIP synergy scores per dose pair", "Predicted vs. observed effect",
                    "Dose-response for each drug alone", "Synergy heatmaps (PDF)"],
        "datasets": ["1 time point", "37 time points × 2 variants (incl. repetitive conditions)"],
        "status": "✅  All datasets complete",
        "sc": TEAL,
    },
    {
        "title": "Genetic Perturbagen",
        "icon":  "🧬",
        "color": GREEN,
        "light": RGBColor(0xEA, 0xF5, 0xEA),
        "what":  "Screens siRNA / shRNA knockdowns to identify genes affecting cell growth.",
        "outputs": ["Normalised confluency per perturbagen vs. control",
                    "Growth rate at each time point", "Multiple control group support"],
        "datasets": ["1 time point, 1 control", "23 time points, multiple controls"],
        "status": "✅  All datasets complete",
        "sc": TEAL,
    },
    {
        "title": "Genetic-Chemical",
        "icon":  "⚗️",
        "color": ORANGE,
        "light": RGBColor(0xFF, 0xF3, 0xE8),
        "what":  "Combines gene knockdown with drug treatment to test gene–drug interactions.",
        "outputs": ["Normalised activity per gene/drug condition",
                    "Synergy scoring between genetic and chemical perturbagens",
                    "Growth rates for combined conditions"],
        "datasets": ["1 time point (53 TP)", "Several time points (53 TP)"],
        "status": "⚠️  Known compatibility issue with example data",
        "sc": ORANGE,
    },
]

positions = [(0.3, 1.05), (6.85, 1.05), (0.3, 4.2), (6.85, 4.2)]

for tp, (lft, top) in zip(TYPES, positions):
    # Card background
    box(s2, lft, top, 6.2, 3.0, fill=tp["light"], line=tp["color"], lw=1.5)
    box(s2, lft, top, 6.2, 0.5, fill=tp["color"])

    # Header
    txt(s2, f"{tp['icon']}  {tp['title']}",
        lft+0.15, top+0.07, 5.2, 0.38, size=16, bold=True, color=WHITE)

    # Status badge
    box(s2, lft+4.3, top+0.09, 1.75, 0.3, fill=tp["sc"])
    txt(s2, tp["status"].split("  ")[0],
        lft+4.32, top+0.1, 1.73, 0.28,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # What it does
    txt(s2, tp["what"], lft+0.15, top+0.57, 5.85, 0.42, size=11, color=TEXT_DARK)

    # Outputs label + list
    txt(s2, "Outputs:", lft+0.15, top+1.05, 1.1, 0.25,
        size=10, bold=True, color=tp["color"])
    bullet_block(s2, tp["outputs"], lft+0.15, top+1.3, 2.85, 1.5, size=10)

    # Datasets label + list
    txt(s2, "Datasets run:", lft+3.2, top+1.05, 2.5, 0.25,
        size=10, bold=True, color=tp["color"])
    bullet_block(s2, tp["datasets"], lft+3.2, top+1.3, 2.9, 1.5, size=10)

footer(s2, 2)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Drug Screen Results
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3, WHITE)
header(s3, "Drug Screen Results — Single Agent Dose-Response",
       "MK-1775, BAY1895344, Prexasertib  ·  MCF7 cells  ·  10K/well  ·  72 h endpoint")

# IC50 table
col_w = [2.1, 1.5, 0.9, 1.0, 2.5]
col_x = [0.3, 2.42, 3.94, 4.86, 5.88]
hdrs  = ["Drug", "IC₅₀ (nM)", "R²", "Fit", "Notes"]

for hd, cx, cw in zip(hdrs, col_x, col_w):
    box(s3, cx, 1.0, cw, 0.42, fill=DARK_BLUE)
    txt(s3, hd, cx+0.06, 1.04, cw-0.1, 0.34,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Prexasertib", "9.98",   "0.967", "Good",
     "Curve plateaus at ~40–50%\n— partial inhibition in tested range",
     True),
    ("BAY1895344",  "272.87", "0.957", "Good",
     "Full inhibition achieved\nat high concentrations",
     False),
    ("MK-1775",     "657.30", "0.989", "Excellent",
     "Cleanest sigmoidal fit;\nhighest R² of the three",
     False),
]
for ri, (drug, ic50, r2, fit, note, hi) in enumerate(rows):
    y = 1.42 + ri * 0.65
    rf = RGBColor(0xE2, 0xEF, 0xDA) if hi else (LIGHT_GREY if ri % 2 else WHITE)
    for val, cx, cw in zip([drug, ic50, r2, fit, note], col_x, col_w):
        box(s3, cx, y, cw, 0.62, fill=rf, line=RGBColor(0xCC,0xCC,0xCC), lw=0.5)
        txt(s3, val, cx+0.07, y+0.06, cw-0.12, 0.52,
            size=11 if val not in [drug, ic50] else (13 if val == ic50 else 12),
            bold=(val == drug),
            color=GREEN if (val == drug and hi) else TEXT_DARK,
            align=PP_ALIGN.LEFT if val in [drug, note] else PP_ALIGN.CENTER)

box(s3, 8.45, 1.42, 1.4, 0.42, fill=GREEN)
txt(s3, "★  Most Potent", 8.47, 1.44, 1.36, 0.38,
    size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Key findings
box(s3, 0.3, 3.45, 8.1, 3.55,
    fill=RGBColor(0xF5, 0xF9, 0xFF), line=MID_BLUE, lw=1)
txt(s3, "Key Findings", 0.5, 3.52, 4.0, 0.35, size=14, bold=True, color=DARK_BLUE)
bullet_block(s3, [
    "Prexasertib is the most potent agent — IC₅₀ ≈ 10 nM, roughly 27× more potent than BAY1895344 and 66× more potent than MK-1775.",
    "Prexasertib's dose-response curve did not reach complete inhibition (plateaus at ~40–50%), suggesting a cytostatic rather than cytotoxic effect at the concentrations tested.",
    "MK-1775 produced the cleanest curve fit (R² = 0.989), indicating a reliable, well-defined sigmoidal dose-response relationship.",
    "Multi-timepoint analysis (37 time points, 74 h) confirmed that potency increases over time for all three drugs, as captured by the growth rate metric.",
    "Single biological replicate — results should be interpreted with caution and validated with additional replicates.",
], 0.5, 3.95, 7.7, 2.8, size=11, marker="▶  ")

# Potency chart (visual bar)
txt(s3, "Relative Potency", 8.6, 3.52, 4.5, 0.32, size=13, bold=True, color=DARK_BLUE)
for i, (drug, ic50_val, color) in enumerate([
    ("Prexasertib", 9.98,   GREEN),
    ("BAY1895344",  272.87, MID_BLUE),
    ("MK-1775",     657.30, DARK_BLUE),
]):
    y = 3.95 + i * 1.05
    bar_w = ic50_val / 657.3 * 4.1
    txt(s3, drug, 8.6, y, 2.2, 0.3, size=11, bold=True, color=color)
    box(s3, 8.6, y+0.33, bar_w, 0.38, fill=color)
    txt(s3, f"IC₅₀ = {ic50_val} nM", 8.6+bar_w+0.08, y+0.36, 1.8, 0.3,
        size=10, color=color, bold=True)

footer(s3, 3)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Drug Combination Synergy
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4, WHITE)
header(s4, "Drug Combination Synergy — Bliss Independence Model",
       "MK-1775 combined with Prexasertib and BAY1895344  ·  MCF7 cells  ·  7×7 dose matrix  ·  72 h")

# Context strip
box(s4, 0.3, 1.0, 12.7, 0.55,
    fill=RGBColor(0xEF, 0xF4, 0xFB), line=MID_BLUE, lw=1)
txt(s4, "Bliss independence assumes the two drugs act independently. "
        "A Bliss score > 0 means the combination is more effective than predicted (synergy); "
        "< 0 means less effective (antagonism); ≈ 0 means additive.",
    0.5, 1.06, 12.3, 0.42, size=11, color=DARK_BLUE)

for label, badge, badge_fill, bg_fill, scores, interpretation, y0 in [
    (
        "MK-1775  +  Prexasertib",
        "SYNERGISTIC", GREEN, RGBColor(0xE8, 0xF5, 0xE9),
        [
            "Max Bliss score: +0.56  (MK-1775 1,235 nM + Prexasertib 4.7 nM)",
            "Strong synergy across mid-range MK-1775 (137–1,235 nM) combined with low-dose Prexasertib (1.4–7.1 nM)",
            "Synergy diminishes at saturating MK-1775 concentrations (≥11,111 nM) — ceiling effect as single agent effect dominates",
            "Pattern consistent across both 1-timepoint and multi-timepoint datasets",
        ],
        "These two checkpoint kinase inhibitors (WEE1 + CHK1/2) show clear synergy at "
        "clinically relevant dose combinations, supporting their combined use in MCF7 cells.",
        1.65,
    ),
    (
        "MK-1775  +  BAY1895344",
        "WEAK / ADDITIVE", ORANGE, RGBColor(0xFF, 0xF8, 0xF0),
        [
            "Max Bliss score: +0.42  (MK-1775 137 nM + BAY1895344 167 nM)",
            "Majority of scores near zero — interaction is largely additive across the matrix",
            "Occasional moderate synergy at specific dose pairs but no consistent pattern",
            "Both 1-timepoint and multi-timepoint datasets show similar additive behaviour",
        ],
        "MK-1775 (WEE1i) and BAY1895344 (ATRi) do not demonstrate robust synergy in MCF7 cells "
        "at the concentrations tested — the combination does not exceed additive expectations.",
        4.4,
    ),
]:
    box(s4, 0.3, y0, 12.7, 2.55, fill=bg_fill, line=badge_fill, lw=1.5)
    box(s4, 0.3, y0, 12.7, 0.48, fill=badge_fill)
    txt(s4, label, 0.5, y0+0.07, 8.0, 0.35, size=15, bold=True, color=WHITE)
    box(s4, 9.7, y0+0.07, 3.1, 0.35, fill=WHITE)
    txt(s4, badge, 9.72, y0+0.08, 3.06, 0.3,
        size=12, bold=True, color=badge_fill, align=PP_ALIGN.CENTER)

    # Bullet points (left)
    bullet_block(s4, scores, 0.45, y0+0.57, 7.5, 1.85, size=11)

    # Interpretation (right)
    box(s4, 8.1, y0+0.55, 4.75, 1.88,
        fill=WHITE, line=badge_fill, lw=0.8)
    txt(s4, "Interpretation", 8.25, y0+0.6, 4.4, 0.28,
        size=10, bold=True, color=badge_fill)
    txt(s4, interpretation, 8.25, y0+0.92, 4.45, 1.4,
        size=10, color=TEXT_DARK, italic=True)

footer(s4, 4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Genetic Perturbagen Screens
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5, WHITE)
header(s5, "Genetic Perturbagen Screen Results",
       "siRNA/shRNA knockdown screens  ·  MCF7 cells  ·  1 and multi-timepoint datasets")

# ── Dataset 1: 1 time point ──────────────────────────────────────────────────
y0 = 1.08
box(s5, 0.3, y0, 12.73, 2.88, fill=RGBColor(0xF5, 0xF9, 0xFF), line=MID_BLUE, lw=1)
box(s5, 0.3, y0, 12.73, 0.52, fill=MID_BLUE)
txt(s5, "1 Time Point  |  1 Control  (48 h endpoint)  —  75 perturbagens screened",
    0.5, y0+0.08, 9.5, 0.36, size=14, bold=True, color=WHITE)
box(s5, 10.5, y0+0.1, 2.35, 0.32, fill=TEAL)
txt(s5, "✅  COMPLETE", 10.52, y0+0.11, 2.31, 0.28,
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s5, "Gene perturbagens tested at 40 ng/well in MCF7 cells (12K/well). "
        "Normalised confluency relative to non-targeting control (100%). "
        "Note: gene names are anonymized in the HTSplotter public example dataset.",
    0.5, y0+0.6, 12.3, 0.32, size=11, italic=True, color=DARK_BLUE)

bullet_block(s5, [
    "75 perturbagens screened in a single endpoint confluency assay",
    "Values above 100% = enhanced growth; below 100% = growth inhibition",
    "Identifies candidate genes whose knockdown significantly alters cell proliferation",
], 0.5, y0+0.98, 5.5, 1.78, size=11)

# Notable hits table (right side of block)
txt(s5, "Notable Hits  (48 h confluency vs. control)", 6.2, y0+0.98, 6.55, 0.28,
    size=11, bold=True, color=DARK_BLUE)

hit_hdrs = ["Target", "Confluency %", "Effect"]
hit_cw   = [1.6, 1.7, 2.85]
hit_cx   = [6.2, 7.82, 9.54]
for hd, cx, cw in zip(hit_hdrs, hit_cx, hit_cw):
    box(s5, cx, y0+1.3, cw, 0.3, fill=DARK_BLUE)
    txt(s5, hd, cx+0.05, y0+1.33, cw-0.08, 0.24,
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

hit_rows = [
    ("GeneBK",  "65.2%",  "↓ Growth inhibition",  RGBColor(0xE2,0xEF,0xDA), GREEN),
    ("GeneBI",  "79.1%",  "↓ Growth inhibition",  WHITE,                    GREEN),
    ("GeneAG",  "79.9%",  "↓ Growth inhibition",  RGBColor(0xE2,0xEF,0xDA), GREEN),
    ("GeneAC", "126.1%",  "↑ Enhanced growth",    RGBColor(0xFF,0xF3,0xE8), ORANGE),
    ("GeneX",  "124.5%",  "↑ Enhanced growth",    WHITE,                    ORANGE),
    ("GeneE",  "124.0%",  "↑ Enhanced growth",    RGBColor(0xFF,0xF3,0xE8), ORANGE),
]
for ri, (gene, pct, effect, rf, ec) in enumerate(hit_rows):
    ry = y0 + 1.6 + ri * 0.34
    for val, cx, cw in zip([gene, pct, effect], hit_cx, hit_cw):
        box(s5, cx, ry, cw, 0.32, fill=rf, line=RGBColor(0xCC,0xCC,0xCC), lw=0.4)
        col = ec if val in [gene, pct] else TEXT_DARK
        txt(s5, val, cx+0.05, ry+0.05, cw-0.08, 0.24,
            size=10, bold=(val == gene), color=col,
            align=PP_ALIGN.CENTER if val == pct else PP_ALIGN.LEFT)

# ── Dataset 2: several time points ───────────────────────────────────────────
y1 = 4.12
box(s5, 0.3, y1, 12.73, 2.88, fill=RGBColor(0xF5, 0xF9, 0xFF), line=MID_BLUE, lw=1)
box(s5, 0.3, y1, 12.73, 0.52, fill=MID_BLUE)
txt(s5, "Several Time Points  |  Multiple Controls  (23 time points, 3 h intervals, ~69 h total)",
    0.5, y1+0.08, 9.5, 0.36, size=14, bold=True, color=WHITE)
box(s5, 10.5, y1+0.1, 2.35, 0.32, fill=TEAL)
txt(s5, "✅  COMPLETE", 10.52, y1+0.11, 2.31, 0.28,
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s5, "Same perturbagens measured across time, enabling growth kinetics and rate calculations per target.",
    0.5, y1+0.6, 12.3, 0.32, size=11, italic=True, color=DARK_BLUE)

bullet_block(s5, [
    "Growth rate computed at each time point as a ratio to matched control",
    "Growth rate > 1.0 — cells proliferating faster (possible tumour suppressor knockdown)",
    "Growth rate < 1.0 — cells proliferating slower (growth-promoting gene silenced)",
    "Multiple independent control groups enable robust plate-wide normalisation",
    "Time-resolved data reveals whether effects are immediate, delayed, or transient",
], 0.5, y1+0.98, 8.5, 1.78, size=11)

box(s5, 9.2, y1+0.95, 3.65, 1.82, fill=WHITE, line=MID_BLUE, lw=0.8)
txt(s5, "Interpretation", 9.35, y1+1.0, 3.4, 0.28, size=10, bold=True, color=MID_BLUE)
txt(s5, "Time-course data provides richer biological insight than endpoint alone — "
        "a perturbagen with a delayed growth rate change may reflect an indirect "
        "or adaptive cellular response.",
    9.35, y1+1.32, 3.4, 1.38, size=10, italic=True, color=TEXT_DARK)

footer(s5, 5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — AI-Assisted Chat Interface
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg(s6, WHITE)
header(s6, "AI-Assisted Results Interpretation",
       "Interactive Claude-powered chat interface built on top of HTSplotter")

# Intro banner
box(s6, 0.3, 1.0, 12.73, 0.72,
    fill=RGBColor(0xEF, 0xF4, 0xFB), line=MID_BLUE, lw=1)
txt(s6,
    "After running any HTSplotter analysis, researchers can ask plain-English questions about "
    "their results using an AI chat interface powered by the Anthropic Claude API (claude-opus-4-7). "
    "No bioinformatics expertise required — just type your question.",
    0.5, 1.07, 12.3, 0.58, size=12, color=DARK_BLUE)

# ── LEFT: how it works ───────────────────────────────────────────────────────
box(s6, 0.3, 1.82, 5.9, 5.23,
    fill=RGBColor(0xF5, 0xF9, 0xFF), line=MID_BLUE, lw=1)
txt(s6, "How It Works", 0.5, 1.9, 5.5, 0.38,
    size=15, bold=True, color=DARK_BLUE)

for i, (step, detail) in enumerate([
    ("Run HTSplotter",
     "Execute any analysis (drug screen, combination, genetic). Results are saved as .txt files in the output_results/ folder."),
    ("Set your API key",
     "Get a free API key from console.anthropic.com. Set it once: export ANTHROPIC_API_KEY='sk-ant-...'"),
    ("Launch the chat",
     "Run:  python chat_results.py  [results_dir]\nThe script loads all result files automatically and caches them for efficiency."),
    ("Ask questions",
     "Type any question about your data. The chat maintains conversation history so you can ask follow-ups."),
    ("Exit anytime",
     "Type quit, exit, or press Ctrl+C to stop."),
]):
    y = 2.38 + i * 0.93
    box(s6, 0.42, y, 0.52, 0.52, fill=MID_BLUE)
    txt(s6, str(i+1), 0.44, y+0.08, 0.48, 0.36,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s6, step, 1.05, y+0.02, 4.95, 0.26, size=11, bold=True, color=DARK_BLUE)
    txt(s6, detail, 1.05, y+0.27, 4.95, 0.6, size=10, color=TEXT_DARK)

# ── RIGHT: example conversation ──────────────────────────────────────────────
box(s6, 6.5, 1.82, 6.53, 5.23,
    fill=RGBColor(0xF0, 0xF4, 0xF9), line=MID_BLUE, lw=1)
txt(s6, "Example Conversation", 6.7, 1.9, 6.1, 0.38,
    size=15, bold=True, color=DARK_BLUE)

def bubble(slide, speaker, text, top, width=6.1, left=6.6):
    user = speaker == "You"
    fc = MID_BLUE if user else RGBColor(0xE8, 0xF1, 0xFB)
    tc = WHITE    if user else TEXT_DARK
    # estimate height based on text length
    lines = text.count("\n") + max(1, len(text) // 62)
    h = 0.3 + lines * 0.2
    box(slide, left, top, width, h,
        fill=fc, line=RGBColor(0xBB, 0xCC, 0xDD), lw=0.5)
    txt(slide, f"{speaker}:  {text}", left+0.12, top+0.06, width-0.18, h-0.1,
        size=10, color=tc)
    return top + h + 0.1

y = 2.4
y = bubble(s6, "You",    "Which drug was most potent?", y)
y = bubble(s6, "Claude", "Prexasertib — IC₅₀ ≈ 10 nM. That's ~27× more potent than BAY1895344 and ~66× more potent than MK-1775 in MCF7 cells.", y)
y = bubble(s6, "You",    "Is there synergy between MK-1775 and Prexasertib?", y)
y = bubble(s6, "Claude", "Yes — Bliss scores reach +0.56, indicating clear synergy. The strongest effect occurs at mid-range MK-1775 (137–1,235 nM) combined with low Prexasertib doses (1.4–7.1 nM).", y)
y = bubble(s6, "You",    "What does a growth rate > 1 mean in the genetic screen?", y)
y = bubble(s6, "Claude", "A growth rate > 1 means those cells grew faster than the control. This could indicate the knockdown silenced a growth-suppressing gene, effectively releasing a brake on proliferation.", y)

# Technical features strip
box(s6, 0.3, 7.08, 12.73, 0.0)
box(s6, 0.3, 6.88, 12.73, 0.22, fill=DARK_BLUE)
feats = [
    "Model: claude-opus-4-7",
    "Streaming responses",
    "Prompt caching (cost-efficient)",
    "Full conversation history",
    "Works with any HTSplotter output folder",
    "github.com/aaronwils246/HTSplotter-chat",
]
txt(s6, "  ·  ".join(feats), 0.4, 6.9, 12.5, 0.18,
    size=9, bold=True, color=LIGHT_BLUE)

footer(s6, 6)


# ── Save ───────────────────────────────────────────────────────────────────
out = "/Users/aaronwilson/HTSplotter/HTSplotter_Summary.pptx"
prs.save(out)
print(f"Saved: {out}")
